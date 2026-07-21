"""
Build: Git Integration & CI/CD — Managing Your Data Pipeline as Code
Snowflake-branded PPTX using the official Jan 2026 template.
"""
import os, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour constants ──────────────────────────────────────────────────────────
DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
VIOLET    = RGBColor(0x7D, 0x44, 0xCF)
PINK      = RGBColor(0xD4, 0x5B, 0x90)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY  = RGBColor(0x71, 0x71, 0x71)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xC8, 0xC8, 0xC8)
GRID_LINE = RGBColor(0xDD, 0xDD, 0xDD)
CONN_LINE = RGBColor(0xCC, 0xCC, 0xCC)

# ── Placeholder helpers ───────────────────────────────────────────────────────
def set_ph(slide, idx, text):
    ph = slide.placeholders[idx]
    t_pos = (ph.top or 0) / 914400
    clean = text.replace('\n', ' ')
    if t_pos < 0.50 and len(clean) > 50:
        print(f"⚠ TITLE TOO LONG: {len(clean)} chars: \"{clean[:50]}...\"")
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{ns}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0')
            pPr.set('marL', '0')

def _pad_body_ph(ph):
    t_pos = (ph.top or 0) / 914400
    if t_pos > 1.20:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
        bodyPr.set('bIns', '91440')

def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame; tf.clear()
    _pad_body_ph(ph)
    lines = [l for l in lines if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size:
            p.font.size = Pt(font_size)

def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame; tf.clear()
    _pad_body_ph(ph)
    first = True
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        if not first:
            pPr = p._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p._p, f'{{{ns}}}pPr')
                p._p.insert(0, pPr)
            spcBef = etree.SubElement(pPr, f'{{{ns}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{ns}}}spcPts')
            spcPts.set('val', '1400')
        first = False
        run = p.add_run(); run.text = heading
        run.font.bold = True; run.font.color.rgb = DK2; run.font.name = "Arial"
        if heading_size: run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph(); bp.level = 1; bp.text = line
            if body_size: bp.font.size = Pt(body_size)

# ── Shape text helper (MANDATORY) ─────────────────────────────────────────────
def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if width <= 2.0 and '\n' not in text and ' ' in text:
        text = text.replace(' ', '\n')
    tf = shape.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.text = text
    p.font.name = "Arial"; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape

def add_code_block(slide, left, top, width, height, lines,
                   bg_colour=None, font_colour=None, font_size=9):
    if bg_colour is None: bg_colour = LIGHT_BG
    if font_colour is None: font_colour = DK1
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_colour
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    text = "\n".join(lines)
    p = tf.paragraphs[0]; p.text = text
    p.font.name = "Arial"; p.font.size = Pt(font_size)
    p.font.color.rgb = font_colour; p.alignment = PP_ALIGN.LEFT
    return shape

def set_table_borders(tbl, n_rows, n_cols):
    for ri in range(n_rows):
        for ci in range(n_cols):
            tc = tbl.cell(ri, ci)._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))
            for edge in ["lnL", "lnR", "lnT", "lnB"]:
                ln = etree.SubElement(tcPr, qn(f"a:{edge}"), w="12700")
                sf = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr"), val="C8C8C8")

# ── Verification ──────────────────────────────────────────────────────────────
from pptx.enum.shapes import MSO_SHAPE_TYPE

def verify_slide(slide, prs, slide_num):
    issues = []
    safe_bottom = 5.10
    for shape in slide.shapes:
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        bot = t + h; right = l + w
        if not shape.is_placeholder and bot > safe_bottom and w > 0.5:
            issues.append(f"  OVERFLOW: shape at ({l:.2f}\",{t:.2f}\") bottom={bot:.2f}\"")
        if not shape.is_placeholder and w > 0.3:
            if right > 9.55:
                issues.append(f"  RIGHT OVERFLOW: shape at ({l:.2f}\",{t:.2f}\") right={right:.2f}\"")
    if issues:
        print(f"⚠ SLIDE {slide_num}:")
        for i in issues: print(i)
    else:
        print(f"✓ SLIDE {slide_num} OK")
    return issues

def verify_deck(prs):
    issues = []
    content_slides = 0; visual_slides = 0
    for slide in prs.slides:
        has_shapes = any(
            not s.is_placeholder and (s.width or 0)/914400 > 0.5 and (s.height or 0)/914400 > 0.3
            for s in slide.shapes
        )
        try:
            layout_name = slide.slide_layout.name
            is_divider = "Quote" in layout_name and "Violet" in layout_name
            is_cover = any(s.is_placeholder and s.placeholder_format.idx == 3 for s in slide.shapes)
            is_thankyou = "Thank" in layout_name
        except:
            is_divider = is_cover = is_thankyou = False
        if not is_cover and not is_divider and not is_thankyou:
            content_slides += 1
            if has_shapes: visual_slides += 1
    ratio = visual_slides / content_slides if content_slides > 0 else 0
    if ratio < 0.40:
        issues.append(f"  LOW VISUAL: {visual_slides}/{content_slides} ({ratio:.0%}) — need ≥40%")
    if issues:
        print(f"⚠ DECK: {issues}")
    else:
        print(f"✓ DECK OK: {len(prs.slides)} slides, {visual_slides}/{content_slides} visual ({ratio:.0%})")
    return issues

# ── Load template ─────────────────────────────────────────────────────────────
TEMPLATE_SEARCH = [
    os.path.join(os.getcwd(), "templates", "snowflake_template.pptx"),
    os.path.expanduser("~/.cortex/skills/900-999_utilities/945-render-pptx/snowflake_template.pptx"),
]
TEMPLATE = next((p for p in TEMPLATE_SEARCH if os.path.isfile(p)), None)
assert TEMPLATE, "snowflake_template.pptx not found"
prs = Presentation(TEMPLATE)

# Remove all sample slides
while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = (sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
           or sldId.get('r:id'))
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

slide_num = 0

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[13])
set_ph(slide, 3, "YOUR PIPELINE\nAS CODE")
set_ph(slide, 0, "Git Integration & CI/CD for Snowflake and Airflow")
set_ph(slide, 2, "Peter Pham  |  Sales Engineer  |  Snowflake  |  2026")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA (visual timeline blocks)
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "SESSION AGENDA")
set_ph(slide, 1, "What we will cover today")

agenda_items = [
    ("01", "The Problem", "Manual deploy\npain points"),
    ("02", "The Approach", "Git as source\nof truth"),
    ("03", "Capabilities", "Snowflake + Git +\nschemachange"),
    ("04", "Pipeline Flow", "End-to-end\nCI/CD walkthrough"),
    ("05", "Architecture", "GitHub + Snowflake\n+ Airflow diagram"),
    ("06", "Live Demo", "Watch it\nin action"),
]
n = len(agenda_items)
box_w = (9.10 - 0.15 * (n - 1)) / n
x = 0.40
colors = [SF_BLUE, DK2, SF_BLUE, DK2, SF_BLUE, DK2]

for i, (num, topic, desc) in enumerate(agenda_items):
    # Number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + box_w/2 - 0.22), Inches(1.35), Inches(0.44), Inches(0.44))
    circ.fill.solid(); circ.fill.fore_color.rgb = colors[i]
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Topic box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.90), Inches(box_w), Inches(0.55))
    box.fill.solid(); box.fill.fore_color.rgb = colors[i]
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = topic
    p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(
        Inches(x), Inches(2.55), Inches(box_w), Inches(0.65))
    d_box.text_frame.word_wrap = True
    p = d_box.text_frame.paragraphs[0]; p.text = desc
    p.font.size = Pt(8); p.font.color.rgb = BODY_GREY
    p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    x += box_w + 0.15

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "THE PROBLEM WITH\nMANUAL DEPLOYS")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SITUATION: Challenge Cards
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "MANUAL DEPLOYS BREAK PRODUCTION")
set_ph(slide, 1, "Three failure modes that Git-driven CI/CD eliminates")

challenges = [
    ("NO AUDIT TRAIL",
     "Who changed that table?\nWhen was this stored procedure\nlast updated? Nobody knows.\nSlack messages are your only log.",
     SF_BLUE, WHITE),
    ("ENVIRONMENT DRIFT",
     "Dev, staging, and prod\ndiverge silently over time.\nA fix in dev never makes it\nto production — until it breaks.",
     DK2, WHITE),
    ("SILENT PIPELINE BREAKS",
     "A DAG change breaks a Snowflake\ndependency with no warning.\nYour data pipeline fails at 2am\nwith no rollback plan.",
     VIOLET, WHITE),
]
card_w = 2.80; gap = 0.25; x = 0.40
for title, body, hdr_color, hdr_txt in challenges:
    # Header bar
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.35), Inches(card_w), Inches(0.45))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = hdr_color
    hdr.line.fill.background()
    tf = hdr.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = hdr_txt; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Body card
    body_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.90), Inches(card_w), Inches(2.80))
    body_shape.fill.solid(); body_shape.fill.fore_color.rgb = LIGHT_BG
    body_shape.line.color.rgb = BORDER; body_shape.line.width = Pt(0.5)
    tf = body_shape.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]; p.text = body
    p.font.size = Pt(10); p.font.color.rgb = DK1; p.font.name = "Arial"

    x += card_w + gap

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — THE APPROACH
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "THE APPROACH")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — APPROACH: 4 Principle Cards (Numbered Steps)
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ONE REPO. ONE SOURCE OF TRUTH.")
set_ph(slide, 1, "Four principles that separate managed pipelines from chaos")

principles = [
    ("1", "EVERYTHING\nIS CODE",
     "Snowflake objects AND\nAirflow DAGs live in Git.\nNothing is created\ndirectly in a UI."),
    ("2", "CHANGES FLOW\nTHROUGH ENVS",
     "dev → staging → prod,\nnever backwards.\nEvery change is\ntested before prod."),
    ("3", "CI/CD\nAUTOMATES DEPLOYS",
     "Lint, test, and deploy\nwithout human error.\nGitHub Actions handles\nevery step."),
    ("4", "SECRETS STAY\nOUT OF GIT",
     "Credentials go in a\nsecrets manager.\nInjected at deploy time,\nnever committed."),
]
n = len(principles); gap = 0.20
col_w = (9.10 - gap * (n - 1)) / n; x = 0.40

for i, (num, title, body) in enumerate(principles):
    color = SF_BLUE if i % 2 == 0 else DK2
    # Number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + col_w/2 - 0.22), Inches(1.35), Inches(0.44), Inches(0.44))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Title box
    t_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.88), Inches(col_w), Inches(0.65))
    t_box.fill.solid(); t_box.fill.fore_color.rgb = color
    t_box.line.fill.background()
    tf = t_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Body box
    b_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.63), Inches(col_w), Inches(2.20))
    b_box.fill.solid(); b_box.fill.fore_color.rgb = LIGHT_BG
    b_box.line.color.rgb = BORDER; b_box.line.width = Pt(0.5)
    tf = b_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(6); tf.margin_bottom = Pt(6)
    p = tf.paragraphs[0]; p.text = body
    p.font.size = Pt(9); p.font.color.rgb = DK1; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    x += col_w + gap

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — CORE CAPABILITIES
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "CORE\nCAPABILITIES")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Snowflake Git Repository (2-col layout)
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "SNOWFLAKE READS FROM GITHUB NATIVELY")
set_ph(slide, 3, "Link any GitHub repo — no custom tooling required")
set_ph_sections(slide, 1, [
    ("How It Works", [
        "CREATE OR REPLACE GIT REPOSITORY command",
        "Links GitHub to Snowflake using a PAT stored as a secret",
        "Pull the latest files from any branch on demand",
        "Compatible with GitHub, GitLab, and Bitbucket",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("What You Can Do", [
        "Execute SQL migration scripts directly from the repo",
        "Deploy stored procedures and UDFs from Git",
        "Sync stage files, configs, and data assets",
        "Trigger deployments from GitHub Actions pipelines",
    ]),
], heading_size=11, body_size=10)
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — schemachange: Migration-Based Versioning
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "SCHEMA CHANGES ARE MIGRATIONS, NOT EDITS")
set_ph(slide, 1, "schemachange tracks every SQL file run — and enforces order")

# Code block showing migration folder
add_code_block(slide, 0.40, 1.35, 4.40, 2.60, [
    "# Migration folder structure:",
    "",
    "/snowflake/",
    "  V1.0.0__create_raw_schema.sql",
    "  V1.1.0__add_customers_table.sql",
    "  V1.2.0__add_orders_view.sql",
    "  V2.0.0__add_region_column.sql",
    "",
    "# schemachange reads CHANGE_HISTORY",
    "# table in Snowflake to track what",
    "# has already been deployed.",
], bg_colour=DK2, font_colour=WHITE, font_size=9)

# Three annotation boxes on the right
annotations = [
    ("VERSIONED FILES",
     "Prefix V + version + description.\nNever edit a deployed file —\nadd a new one instead.",
     SF_BLUE, WHITE),
    ("APPEND-ONLY RULE",
     "schemachange rejects out-of-order\nfiles. This prevents accidental\nrollbacks in production.",
     DK2, WHITE),
    ("GITHUB ACTIONS TRIGGER",
     "On merge to main, CI/CD runs\nschemachange automatically.\nZero manual steps to deploy.",
     BODY_GREY, WHITE),
]
y = 1.35
for title, body, color, txt in annotations:
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.00), Inches(y), Inches(4.50), Inches(0.35))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = txt; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    body_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.00), Inches(y + 0.40), Inches(4.50), Inches(0.65))
    body_box.fill.solid(); body_box.fill.fore_color.rgb = LIGHT_BG
    body_box.line.color.rgb = BORDER; body_box.line.width = Pt(0.5)
    tf = body_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]; p.text = body
    p.font.size = Pt(9); p.font.color.rgb = DK1; p.font.name = "Arial"
    y += 1.15

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Monorepo Structure (folder tree visual)
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ONE REPO GOVERNS BOTH SIDES OF YOUR PIPELINE")
set_ph(slide, 1, "Snowflake migrations and Airflow DAGs co-deploy from a single source")

# Left side: folder tree as code block
add_code_block(slide, 0.40, 1.35, 4.20, 3.60, [
    "my-data-pipeline/",
    "├── snowflake/",
    "│   ├── V1.0.0__init_schema.sql",
    "│   ├── V1.1.0__add_customers.sql",
    "│   └── V2.0.0__new_column.sql",
    "├── dags/",
    "│   ├── ingest_customers_dag.py",
    "│   └── daily_reporting_dag.py",
    "├── tests/",
    "│   └── validate_migrations.sql",
    "└── .github/",
    "    └── workflows/",
    "        └── deploy.yml",
], bg_colour=DK2, font_colour=WHITE, font_size=9)

# Right side: 4 callout boxes explaining each folder
folders = [
    ("/snowflake/", "SQL migration files in version order.\nschemachange deploys new\nfiles to Snowflake on merge.", SF_BLUE, WHITE),
    ("/dags/", "Airflow DAG Python files.\nCI/CD syncs these to S3,\nMWAA, or Astronomer.", DK2, WHITE),
    ("/tests/", "SQL and Python validation\nscripts. CI runs these on\nevery Pull Request.", BODY_GREY, WHITE),
    ("/.github/workflows/", "GitHub Actions YAML. Defines\nthe CI/CD pipeline that\nruns all deployments.", SF_BLUE, WHITE),
]
y = 1.35
for title, body, color, txt in folders:
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   4.80, y, 4.68, 0.33, title, color, txt,
                   font_size=9, bold=True)
    body_s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.80), Inches(y + 0.37), Inches(4.68), Inches(0.42))
    body_s.fill.solid(); body_s.fill.fore_color.rgb = LIGHT_BG
    body_s.line.color.rgb = BORDER; body_s.line.width = Pt(0.5)
    tf = body_s.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]; p.text = body
    p.font.size = Pt(8); p.font.color.rgb = DK1; p.font.name = "Arial"
    y += 0.88

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — END-TO-END PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[20])
set_ph(slide, 1, "END-TO-END\nCOMMIT TO PRODUCTION")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CI/CD Flow: 5-Step Chevron
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "FROM COMMIT TO PRODUCTION IN 5 STEPS")
set_ph(slide, 1, "A single Pull Request deploys both Snowflake and Airflow changes")

stages = [
    ("1", "OPEN\nPULL REQUEST", "Developer commits\nchanges to a feature\nbranch in GitHub"),
    ("2", "CI VALIDATES", "GitHub Actions lints\nSQL, validates DAG\nsyntax, runs tests"),
    ("3", "PR REVIEW\n& MERGE", "Team reviews and\napproves. Merge to\nmain triggers deploy"),
    ("4", "CD DEPLOYS", "schemachange runs\nnew migrations on\nSnowflake automatically"),
    ("5", "AIRFLOW SYNCS", "DAG files synced\nto Airflow S3 bucket.\nNew DAGs active"),
]
n = len(stages)
chev_w = 1.85; overlap = 0.18
total_w = chev_w * n - overlap * (n - 1)
x_start = (10.0 - total_w) / 2

for i, (num, label, desc) in enumerate(stages):
    x = x_start + i * (chev_w - overlap)
    color = SF_BLUE if i % 2 == 0 else DK2

    # Number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + chev_w/2 - 0.18), Inches(1.28), Inches(0.36), Inches(0.36))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Chevron
    chev = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x), Inches(1.70), Inches(chev_w), Inches(0.65))
    chev.fill.solid(); chev.fill.fore_color.rgb = color
    chev.line.fill.background()
    tf = chev.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(8); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Description
    d = slide.shapes.add_textbox(
        Inches(x + 0.10), Inches(2.45), Inches(chev_w - 0.20), Inches(1.00))
    d.text_frame.word_wrap = True
    p = d.text_frame.paragraphs[0]; p.text = desc
    p.font.size = Pt(8); p.font.color.rgb = DK1
    p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

# Bottom banner
banner = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.70), Inches(9.10), Inches(0.55))
banner.fill.solid(); banner.fill.fore_color.rgb = DK2
banner.line.fill.background()
tf = banner.text_frame; tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "This same flow works for EVERY change — a new table, a new DAG, a new stored procedure"
p.font.size = Pt(10); p.font.bold = True
p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "CI/CD VS. MANUAL")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "CI/CD VS. MANUAL: EVERY STEP MATTERS")
set_ph(slide, 1, "Each manual step is a potential failure point — CI/CD eliminates them")

headers = ["Capability", "Manual Approach", "Snowflake + Git CI/CD"]
data = [
    ["Version history", "Slack messages + tribal knowledge", "Git commit log with author, timestamp, diff"],
    ["Environment promotion", "Copy-paste SQL between environments", "Automated PR flow: dev → staging → prod"],
    ["Schema migrations", "ALTER TABLE run directly in prod", "schemachange: versioned, ordered, tracked"],
    ["Airflow DAG deploy", "SCP/SFTP files to Airflow server", "CI/CD syncs DAG files to S3/MWAA on merge"],
    ["Secrets management", "Credentials hardcoded in scripts", "GitHub Secrets injected at runtime only"],
    ["Audit trail", "None — no record of what changed", "Every change tied to PR, author, timestamp"],
    ["Rollback", "Manual re-run of previous SQL", "Revert the PR, re-run schemachange"],
]
n_rows = len(data) + 1; n_cols = 3
tbl_shape = slide.shapes.add_table(
    n_rows, n_cols,
    Inches(0.40), Inches(1.30), Inches(9.10), Inches(0.47 * n_rows))
tbl = tbl_shape.table

col_widths = [1.80, 3.40, 3.90]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = Inches(cw)

for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DK2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci); cell.text = val
        cell.fill.solid()
        if ci == 0:
            cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.bold = True
                p.font.color.rgb = DK1; p.font.name = "Arial"
        elif ci == 1:
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else RGBColor(0xFF,0xEE,0xEE)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = BODY_GREY
                p.font.name = "Arial"
        else:
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = DK1
                p.font.name = "Arial"

set_table_borders(tbl, n_rows, n_cols)
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[22])
set_ph(slide, 1, "THE ARCHITECTURE")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Architecture: Layered Stack
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "GITHUB IS THE CONTROL PLANE")
set_ph(slide, 1, "One repo, one pipeline, two deployment targets")

# Layer 1: GitHub (top)
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
               0.40, 1.35, 9.10, 0.55,
               "GITHUB  |  Migrations + DAGs + Tests + Workflows — single source of truth",
               DK2, WHITE, font_size=11, bold=True)

# Layer 2: GitHub Actions
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
               0.40, 2.00, 9.10, 0.55,
               "GITHUB ACTIONS  |  On PR: lint SQL, validate DAGs, run tests     On merge: deploy to targets",
               SF_BLUE, WHITE, font_size=11, bold=True)

# Layer 3 label
lbl = slide.shapes.add_textbox(Inches(0.40), Inches(2.68), Inches(9.10), Inches(0.22))
p = lbl.text_frame.paragraphs[0]; p.text = "DEPLOYMENT TARGETS"
p.font.size = Pt(9); p.font.bold = True
p.font.color.rgb = BODY_GREY; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

# Left target: Snowflake
snow_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(2.97), Inches(4.35), Inches(1.80))
snow_box.fill.solid(); snow_box.fill.fore_color.rgb = LIGHT_BG
snow_box.line.color.rgb = SF_BLUE; snow_box.line.width = Pt(1.5)
tf = snow_box.text_frame; tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(6)

# Snowflake header inside box
p = tf.paragraphs[0]
r = p.add_run(); r.text = "SNOWFLAKE"
r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = SF_BLUE; r.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER

for bullet in [
    "• Git Repository object links GitHub repo",
    "• schemachange applies new SQL migrations",
    "• Tables, views, stored procs deployed as code",
    "• CHANGE_HISTORY table tracks every deploy",
]:
    bp = tf.add_paragraph(); bp.text = bullet
    bp.font.size = Pt(9); bp.font.color.rgb = DK1; bp.font.name = "Arial"

# Right target: Airflow
af_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.15), Inches(2.97), Inches(4.35), Inches(1.80))
af_box.fill.solid(); af_box.fill.fore_color.rgb = LIGHT_BG
af_box.line.color.rgb = DK2; af_box.line.width = Pt(1.5)
tf = af_box.text_frame; tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(6)

p = tf.paragraphs[0]
r = p.add_run(); r.text = "AIRFLOW"
r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = DK2; r.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER

for bullet in [
    "• DAG files (Python) live in /dags/ folder",
    "• CI/CD syncs changed files to S3 bucket",
    "• MWAA, Astronomer, or self-hosted",
    "• DAGs call Snowflake via Airflow Snowflake hook",
]:
    bp = tf.add_paragraph(); bp.text = bullet
    bp.font.size = Pt(9); bp.font.color.rgb = DK1; bp.font.name = "Arial"

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# CHAPTER DIVIDER — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "LIVE DEMO")
verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Demo Walkthrough: Numbered Steps
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "DEMO: END-TO-END CI/CD IN ACTION")
set_ph(slide, 1, "Watch a schema change and DAG update deploy from one Pull Request")

demo_steps = [
    ("1", "Show the repo structure in GitHub",
     "/snowflake/, /dags/, /.github/workflows/deploy.yml"),
    ("2", "Create a new SQL migration file",
     "V2.0.0__add_email_column.sql — adds a column to the CUSTOMERS table"),
    ("3", "Open a Pull Request",
     "CI runs automatically: lint SQL, validate syntax, unit tests pass"),
    ("4", "Merge the PR",
     "GitHub Actions triggers deploy.yml — schemachange detects new file"),
    ("5", "Verify in Snowflake",
     "DESCRIBE TABLE CUSTOMERS — new EMAIL column is live in production"),
    ("6", "DAG change deploys simultaneously",
     "Updated DAG is live in Airflow within 60 seconds of merge"),
]

# 2 columns of 3 steps each
for i, (num, title, detail) in enumerate(demo_steps):
    col = i % 2; row = i // 2
    x = 0.40 + col * 4.65
    y = 1.35 + row * 1.18
    color = SF_BLUE if col == 0 else DK2

    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    step_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.50), Inches(y), Inches(3.85), Inches(0.40))
    step_box.fill.solid(); step_box.fill.fore_color.rgb = color
    step_box.line.fill.background()
    tf = step_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"

    det_box = slide.shapes.add_textbox(
        Inches(x + 0.50), Inches(y + 0.45), Inches(3.85), Inches(0.60))
    det_box.text_frame.word_wrap = True
    p = det_box.text_frame.paragraphs[0]; p.text = detail
    p.font.size = Pt(8); p.font.color.rgb = BODY_GREY; p.font.name = "Arial"

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CTA: Three Steps to Get Started
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "THREE STEPS TO GET STARTED THIS WEEK")
set_ph(slide, 1, "Concrete actions to take after this session")

cta_steps = [
    ("STEP 1",
     "CREATE A SNOWFLAKE\nGIT REPOSITORY",
     "Run CREATE OR REPLACE GIT REPOSITORY\nin your Snowflake account and link\nit to your GitHub repo with a PAT.\nVerify with SHOW GIT REPOSITORIES.",
     SF_BLUE, WHITE),
    ("STEP 2",
     "SET UP YOUR FIRST\nMIGRATION FOLDER",
     "Create /snowflake/ in your repo.\nAdd V1.0.0__init.sql with your\nexisting schema DDL. Run schemachange\nlocally to test the connection.",
     DK2, WHITE),
    ("STEP 3",
     "ADD A GITHUB ACTIONS\nWORKFLOW",
     "Create .github/workflows/deploy.yml.\nTrigger on push to main. Add steps\nfor schemachange and DAG sync.\nTest with a small PR.",
     VIOLET, WHITE),
]

n = len(cta_steps); gap = 0.25
col_w = (9.10 - gap * (n - 1)) / n; x = 0.40

for step_label, title, body, color, txt in cta_steps:
    # Step label
    lbl_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.35), Inches(col_w), Inches(0.30))
    lbl_box.fill.solid(); lbl_box.fill.fore_color.rgb = BODY_GREY
    lbl_box.line.fill.background()
    tf = lbl_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = step_label
    p.font.size = Pt(8); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Title box
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.73), Inches(col_w), Inches(0.70))
    title_box.fill.solid(); title_box.fill.fore_color.rgb = color
    title_box.line.fill.background()
    tf = title_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = txt; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER

    # Body box
    body_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.52), Inches(col_w), Inches(2.30))
    body_box.fill.solid(); body_box.fill.fore_color.rgb = LIGHT_BG
    body_box.line.color.rgb = BORDER; body_box.line.width = Pt(0.5)
    tf = body_box.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]; p.text = body
    p.font.size = Pt(9); p.font.color.rgb = DK1; p.font.name = "Arial"

    x += col_w + gap

verify_slide(slide, prs, slide_num)

# ═══════════════════════════════════════════════════════════════════════════════
# THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
slide_num += 1
slide = prs.slides.add_slide(prs.slide_layouts[28])
set_ph(slide, 1, "THANK\nYOU")
verify_slide(slide, prs, slide_num)

# ── Final verification & save ─────────────────────────────────────────────────
verify_deck(prs)

output_path = "/Users/peterpham/projects/Enablement/Git_CICD/Git_CICD_Presentation.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"\n✅ Saved: {output_path}")
print(f"   {slide_num} slides total")
